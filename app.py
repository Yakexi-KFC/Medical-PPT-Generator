import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import io
import json
import base64
import requests
from openai import OpenAI
from PIL import Image  # 用于处理超大图片压缩

st.set_page_config(page_title="肿瘤病例 PPT 生成工作台", layout="wide")

# ==========================================
# 🔑 密钥配置区 (使用 Streamlit Secrets 保护)
# ==========================================
BAIDU_API_KEY = st.secrets["BAIDU_API_KEY"]
BAIDU_SECRET_KEY = st.secrets["BAIDU_SECRET_KEY"]
DEEPSEEK_API_KEY = st.secrets["DEEPSEEK_API_KEY"]

# ==========================================
# 1. 百度 OCR 图片识别模块 (包含超大图防崩溃压缩)
# ==========================================
def get_baidu_access_token():
    url = f"https://aip.baidubce.com/oauth/2.0/token?grant_type=client_credentials&client_id={BAIDU_API_KEY}&client_secret={BAIDU_SECRET_KEY}"
    headers = {'Content-Type': 'application/json', 'Accept': 'application/json'}
    response = requests.request("POST", url, headers=headers, data="")
    return response.json().get("access_token")

def perform_ocr(image_bytes, access_token):
    try:
        # 基础防崩溃压缩：仅当图片真的大于 3.5MB 时，才做轻微的体积压缩
        if len(image_bytes) > 3.5 * 1024 * 1024:
            img = Image.open(io.BytesIO(image_bytes))
            if img.mode != 'RGB':
                img = img.convert('RGB')
            output = io.BytesIO()
            # 仅降低一点保存质量，不改变长宽，防止摩尔纹扭曲
            img.save(output, format="JPEG", quality=70) 
            image_bytes = output.getvalue()

        url = "https://aip.baidubce.com/rest/2.0/ocr/v1/accurate_basic?access_token=" + access_token
        img_base64 = base64.b64encode(image_bytes).decode('utf-8')
        payload = {'image': img_base64}
        headers = {'Content-Type': 'application/x-www-form-urlencoded', 'Accept': 'application/json'}
        response = requests.request("POST", url, headers=headers, data=payload)
        result_json = response.json()
        
        if "words_result" in result_json:
            text_list = [item["words"] for item in result_json["words_result"]]
            return "\n".join(text_list)
        else:
            return f"[识别错误: {result_json.get('error_msg', '未知错误')}]"
    except Exception as e:
        return f"[请求异常: {str(e)}]"

# ==========================================
# 2. AI 结构化提取模块 (学术级深度总结 + 严谨分线)
# ==========================================
def extract_complex_case(patient_text):
    client = OpenAI(
        api_key=DEEPSEEK_API_KEY, 
        base_url="https://api.deepseek.com"
    )
    
    # 【优化核心】：放权临床推理，锁死输出接口
    system_prompt = """
    你是一位顶级的肿瘤内科专家，正在梳理一份复杂的临床病历，准备进行高水平的学术会议汇报（如胃肠肿瘤或妇科肿瘤领域的病例探讨）。
    
    【核心任务与自由度】
    1. 自由梳理逻辑：请发挥你的专业临床判断力，自主分析患者的疾病进展时间轴。你来决定如何划分治疗线数（一线、二线、维持治疗等），并准确判断不同阶段的疗效转归（PR/SD/PD 等）。
    2. 深度医学提炼：不要单纯当一个“文字搬运工”。请计算关键生存指标，评估治疗策略的得失，敏锐捕捉病程中的矛盾点或亮点（例如：特定靶向药跨线使用的疗效、某种耐药机制的出现等）。
    
    【系统接口规范（极度重要）】
    为了对接下游的 PPT 自动渲染系统，你**必须且只能**输出一个标准的 JSON 对象。
    严禁改变以下任何一个键名（Key），你可以根据你的临床推理自由填充对应的值（Value）：
    
    ```json
    {
        "cover": {"title": "晚期XXX癌综合治疗病例汇报"},
        "baseline": {
            "patient_info": "患者姓名(姓氏)、性别、年龄",
            "chief_complaint": "主诉",
            "diagnosis": "完整的临床及病理诊断",
            "key_exams": "关键基线检查"
        },
        "treatments": [
            {
                "phase": "阶段名称（由你自主判断，如：一线治疗 / 维持治疗）", 
                "duration": "具体时间段", 
                "regimen": "完整的用药方案及局部治疗手段", 
                "imaging": "影像学评估结果",
                "markers": "肿瘤标志物变化"
            }
        ],
        "current_admission": {
            "exams": ["检验异常指标1", "检验异常指标2"],
            "imaging": "本次核心影像结论",
            "plan": ["后续治疗计划或考量1", "考量2"]
        },
        "timeline_events": [
            {
                "date": "年月", 
                "phase": "线数或阶段",
                "event_type": "Treatment 或 Evaluation",
                "event": "高度凝练的事件短语"
            }
        ],
        "summary": {
            "highlights": [
                "由你提炼的病例亮点1", 
                "由你提炼的病例亮点2"
            ],
            "discussion": [
                "值得探讨的临床深度问题1",
                "值得探讨的临床深度问题2"
            ]
        }
    }
    ```
    """
    
    response = client.chat.completions.create(
        model="deepseek-v4-pro",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": patient_text}
        ]
        # 注意：移除了 response_format，因为 reasoner 模型不支持强制 JSON 模式
    )
    
    # 获取模型的最终输出内容（忽略前面冗长的 <think> 推理过程）
    raw_content = response.choices[0].message.content
    
    # 增加鲁棒性清洗：确保去除 Markdown 的代码块标记，提取纯 JSON 字符串
    try:
        # 去除可能包含的 ```json 和 ``` 标记
        if "```json" in raw_content:
            json_str = raw_content.split("```json")[1].split("```")[0].strip()
        elif "```" in raw_content:
            json_str = raw_content.split("```")[1].split("```")[0].strip()
        else:
            json_str = raw_content.strip()
            
        return json.loads(json_str)
        
    except json.JSONDecodeError as e:
        # 如果模型偶尔没有严格遵守 JSON 格式，返回友好的报错信息
        raise ValueError(f"AI 生成的数据无法解析为 JSON，请重试。原始返回摘要：{raw_content[:100]}...")

# ==========================================
# 3. 网页端 Markdown 逻辑流生成器 (备用Cheat Sheet)
# ==========================================
def render_logic_line_markdown(data):
    """将 JSON 转化为一目了然的 Markdown 病例逻辑流"""
    lines = []
    
    # 1. 基线部分
    base = data.get("baseline", {})
    lines.append(f"#### 👤 {base.get('patient_info', '患者')} | {base.get('diagnosis', '未提供诊断')}")
    lines.append(f"> **关键基线检查**：{base.get('key_exams', '无')}\n")
    
    # 2. 治疗演变 (垂直流)
    for tx in data.get("treatments", []):
        lines.append(f"**⬇️ {tx.get('phase', '阶段治疗')}** `({tx.get('duration', '')})`")
        lines.append(f"- **方案**：{tx.get('regimen', '')}")
        lines.append(f"- **评估**：{tx.get('imaging', '')} | {tx.get('markers', '')}\n")
        
    # 3. 转归部分
    adm = data.get("current_admission")
    if adm:
        lines.append(f"**➡️ 本次转归与计划**")
        exams = "；".join(adm.get("exams", [])) if isinstance(adm.get("exams", []), list) else str(adm.get("exams", ""))
        lines.append(f"- **异常指标**：{exams}")
        lines.append(f"- **影像评估**：{adm.get('imaging', '')}")
        plan = "；".join(adm.get("plan", [])) if isinstance(adm.get("plan", []), list) else str(adm.get("plan", ""))
        lines.append(f"- **后续处理**：{plan}")
        
    return "\n".join(lines)

# ==========================================
# 4. PPT 生成模块
# ==========================================
class AdvancedPPTMaker:
    def __init__(self, data):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333) 
        self.prs.slide_height = Inches(7.5)
        self.data = self.clean_data(data)
        self.C_PRI = RGBColor(115, 21, 40)   
        self.C_ACC = RGBColor(0, 51, 102)  

    def clean_data(self, data):
        has_surgery = False
        full_text = json.dumps(data, ensure_ascii=False)
        if "根治术" in full_text or "切除术" in full_text or "手术切除" in full_text:
            has_surgery = True
            
        if not has_surgery:
            for tx in data.get("treatments", []):
                if "辅助" in tx.get("phase", ""):
                    tx["phase"] = "一线治疗" 
            for evt in data.get("timeline_events", []):
                if "辅助" in evt.get("phase", ""):
                    evt["phase"] = "一线"
        return data

    def add_header(self, slide, text):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.C_PRI
        shape.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.05), Inches(10), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    def make_cover(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.C_PRI
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10), Inches(2))
        p = tb.text_frame.paragraphs[0]
        p.text = self.data.get("cover", {}).get("title", "病例汇报")
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    def make_baseline(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "病例介绍 (基线资料)")
        base_data = self.data.get("baseline", {})
        content = f"【患者信息】 {base_data.get('patient_info', '')}\n\n" \
                  f"【主诉】 {base_data.get('chief_complaint', '')}\n\n" \
                  f"【临床诊断】\n{base_data.get('diagnosis', '')}\n\n" \
                  f"【关键检查/病理】\n{base_data.get('key_exams', '')}"
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(20) 
        
    def make_treatments(self):
        for tx in self.data.get("treatments", []):
            phase_name = tx.get('phase', '阶段治疗')
            if "辅助" in phase_name and len(tx.get('regimen', '')) < 5:
                continue
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self.add_header(slide, f"治疗经过：{phase_name}")
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(6))
            tf = tb.text_frame
            tf.word_wrap = True 
            p1 = tf.paragraphs[0]
            p1.text = f"【治疗时间】 {tx.get('duration', '')}"
            p1.font.size = Pt(20) 
            p1.font.bold = True
            p1.font.color.rgb = self.C_PRI
            p2 = tf.add_paragraph()
            p2.text = f"\n【用药方案及局部治疗】\n{tx.get('regimen', '')}"
            p2.font.size = Pt(16) 
            p3 = tf.add_paragraph()
            p3.text = f"\n【影像学评估】\n{tx.get('imaging', '')}"
            p3.font.size = Pt(16) 
            p3.font.color.rgb = RGBColor(50, 50, 50)
            p4 = tf.add_paragraph()
            p4.text = f"\n【肿瘤标志物】\n{tx.get('markers', '')}"
            p4.font.size = Pt(16) 
            p4.font.color.rgb = self.C_ACC

    def make_current_admission(self):
        adm_data = self.data.get("current_admission")
        if not adm_data: return
        exams_list = adm_data.get("exams", [])
        exams_str = "\n".join([f"• {item}" for item in exams_list]) if isinstance(exams_list, list) else str(exams_list)
        imaging_str = adm_data.get("imaging", "")
        plan_list = adm_data.get("plan", [])
        plan_str = "\n".join([f"• {item}" for item in plan_list]) if isinstance(plan_list, list) else str(plan_list)
        total_len = len(exams_str) + len(imaging_str) + len(plan_str)
        is_split = len(plan_str) > 200 or total_len > 500
        
        if is_split:
            slide1 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self.add_header(slide1, "本次入院评估 (1/2)")
            tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))
            tf1 = tb1.text_frame
            p_ex_title = tf1.paragraphs[0]
            p_ex_title.text = "【入院检验指标】"
            p_ex_title.font.bold = True
            p_ex_title.font.size = Pt(20)
            p_ex_title.font.color.rgb = self.C_PRI
            p_ex_body = tf1.add_paragraph()
            p_ex_body.text = exams_str + "\n"
            p_ex_body.font.size = Pt(18)
            p_im_title = tf1.add_paragraph()
            p_im_title.text = "【影像学评估】"
            p_im_title.font.bold = True
            p_im_title.font.size = Pt(20)
            p_im_title.font.color.rgb = self.C_PRI
            p_im_body = tf1.add_paragraph()
            p_im_body.text = imaging_str
            p_im_body.font.size = Pt(18)
            
            slide2 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self.add_header(slide2, "后续治疗与随访计划 (2/2)")
            tb2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))
            tf2 = tb2.text_frame
            p_pl_title = tf2.paragraphs[0]
            p_pl_title.text = "【治疗与随访计划】"
            p_pl_title.font.bold = True
            p_pl_title.font.size = Pt(20)
            p_pl_title.font.color.rgb = self.C_PRI
            p_pl_body = tf2.add_paragraph()
            p_pl_body.text = plan_str
            p_pl_body.font.size = Pt(18)
        else:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self.add_header(slide, "本次入院评估及计划 (转归)")
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(6))
            tf = tb.text_frame
            content = f"【入院检验指标】\n{exams_str}\n\n【影像学评估】\n{imaging_str}\n\n【后续计划】\n{plan_str}"
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(16)

    def make_timeline(self):
        events = self.data.get("timeline_events", [])
        if not events: return
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "全病程时间轴概览 (Timeline)")
        line_y = Inches(4.2)
        start_x = Inches(0.6)
        total_width = 12.1 
        count = min(len(events), 12)
        if count > 9:
            card_width = Inches(0.95); card_height = Inches(1.4); font_size_date = Pt(9); font_size_body = Pt(8)
        elif count > 6:
            card_width = Inches(1.3); card_height = Inches(1.2); font_size_date = Pt(10); font_size_body = Pt(9)
        else:
            card_width = Inches(1.6); card_height = Inches(1.2); font_size_date = Pt(12); font_size_body = Pt(11)

        main_line = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, start_x - Inches(0.2), line_y - Inches(0.05), Inches(total_width + 0.4), Inches(0.1))
        main_line.fill.solid()
        main_line.fill.fore_color.rgb = RGBColor(220, 220, 220) 
        main_line.line.fill.background()
        
        for i, evt in enumerate(events[:12]): 
            if count > 1: x = start_x + Inches(total_width * (i / (count - 1)))
            else: x = start_x + Inches(total_width / 2)
            event_text = evt.get("event", "")
            phase_text = evt.get("phase", "") 
            event_type = evt.get("event_type", "Treatment")
            is_pd = "进展" in event_text or "PD" in event_text.upper() or "复发" in event_text
            is_control = "PR" in event_text.upper() or "SD" in event_text.upper() or "缩小" in event_text
            if is_pd: node_color = RGBColor(220, 50, 50) 
            elif is_control and event_type == "Evaluation": node_color = RGBColor(46, 139, 87) 
            else: node_color = self.C_PRI 
            
            stem_height = Inches(1.0)
            stem_top = line_y - stem_height if i % 2 == 0 else line_y
            stem = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, stem_top, Inches(0.03), stem_height) 
            stem.fill.solid()
            stem.fill.fore_color.rgb = node_color
            stem.line.fill.background()
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.15), line_y - Inches(0.15), Inches(0.3), Inches(0.3))
            circle.fill.solid()
            circle.fill.fore_color.rgb = node_color
            circle.line.color.rgb = RGBColor(255, 255, 255); circle.line.width = Pt(2)
            
            card_top = line_y - stem_height - card_height if i % 2 == 0 else line_y + stem_height
            card_x = x - (card_width / 2)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_top, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(250, 250, 250) 
            card.line.color.rgb = node_color; card.line.width = Pt(1.5)
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05); tf.margin_top = Inches(0.05)
            p0 = tf.paragraphs[0]
            p0.text = evt.get("date", "")
            p0.font.bold = True; p0.font.size = font_size_date; p0.font.color.rgb = node_color; p0.alignment = PP_ALIGN.CENTER
            
            if phase_text and phase_text != "评估":
                p_phase = tf.add_paragraph()
                p_phase.text = f"【{phase_text}】"
                p_phase.font.size = font_size_body; p_phase.font.bold = True; p_phase.font.color.rgb = node_color; p_phase.alignment = PP_ALIGN.CENTER
            elif event_type == "Evaluation":
                p_phase = tf.add_paragraph()
                p_phase.text = "【疗效评估】"
                p_phase.font.size = font_size_body; p_phase.font.bold = True; p_phase.font.color.rgb = node_color; p_phase.alignment = PP_ALIGN.CENTER
            p1 = tf.add_paragraph()
            p1.text = event_text
            p1.font.size = font_size_body; p1.font.color.rgb = RGBColor(30, 30, 30); p1.alignment = PP_ALIGN.CENTER

    def make_summary(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_header(slide, "病例思考与总结")
        summary_data = self.data.get("summary", {})
        highlights = []
        discussion = []
        if isinstance(summary_data, list):
            highlights = summary_data
        elif isinstance(summary_data, dict):
            highlights = summary_data.get("highlights", [])
            discussion = summary_data.get("discussion", [])

        top_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.0))
        tf_top = top_box.text_frame
        tf_top.word_wrap = True
        
        for item in highlights:
            p = tf_top.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(22) 
            p.font.bold = True
            p.space_after = Pt(18)
            
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = self.C_PRI 
        line.line.fill.background()

        if discussion:
            bottom_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.8))
            tf_bottom = bottom_box.text_frame
            tf_bottom.word_wrap = True
            p_title = tf_bottom.paragraphs[0]
            p_title.text = "思考："
            p_title.font.size = Pt(22)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(0, 0, 0)
            p_title.space_after = Pt(12)
            
            for item in discussion:
                p = tf_bottom.add_paragraph()
                p.text = f"➤ {item}" 
                p.font.size = Pt(20)
                p.font.bold = True
                p.space_after = Pt(14)

    def build(self):
        self.make_cover()
        self.make_baseline()
        self.make_treatments()
        self.make_current_admission()
        self.make_timeline()
        self.make_summary()
        ppt_stream = io.BytesIO()
        self.prs.save(ppt_stream)
        ppt_stream.seek(0)
        return ppt_stream

def inject_frontend_style():
    st.markdown(
        """
        <style>
        :root {
            --primary: #b00086;
            --primary-dark: #7e0067;
            --accent: #ec008c;
            --ink: #25232a;
            --muted: #6f6b76;
            --line: #ded9e1;
            --panel: #ffffff;
            --soft: #f7f5f8;
            --warning-bg: #fff4fa;
            --warning-line: #f3a4d1;
            --shadow: 0 18px 45px rgba(70, 58, 76, 0.10);
        }

        .stApp {
            background: linear-gradient(180deg, #faf8fb 0%, #f0edf2 44%, #fbfafc 100%);
            color: var(--ink);
        }

        .block-container {
            max-width: 1180px;
            padding-top: 1.5rem;
            padding-bottom: 3rem;
        }

        h1, h2, h3, p, label, span {
            letter-spacing: 0;
        }

        div[data-testid="stTabs"] button {
            min-height: 48px;
            border-radius: 8px 8px 0 0;
            font-weight: 800;
        }

        div[data-testid="stTabs"] button[aria-selected="true"] {
            color: var(--primary);
            border-color: var(--line);
            background: #ffffff;
        }

        .stButton > button,
        .stDownloadButton > button {
            width: 100%;
            min-height: 46px;
            border: 1px solid var(--primary);
            border-radius: 8px;
            background: var(--primary);
            color: #ffffff;
            font-weight: 850;
            box-shadow: 0 10px 20px rgba(176, 0, 134, 0.16);
        }

        .stButton > button:hover,
        .stDownloadButton > button:hover {
            border-color: var(--primary-dark);
            background: var(--primary-dark);
            color: #ffffff;
        }

        .stTextArea textarea,
        .stFileUploader section {
            border-radius: 8px;
        }

        .case-hero {
            position: relative;
            overflow: hidden;
            border: 1px solid rgba(122, 31, 54, 0.14);
            border-radius: 8px;
            padding: 34px 318px 30px 38px;
            background:
                linear-gradient(135deg, rgba(176, 0, 134, 0.95), rgba(236, 0, 140, 0.88)),
                repeating-linear-gradient(90deg, rgba(255,255,255,0.10) 0, rgba(255,255,255,0.10) 1px, transparent 1px, transparent 22px);
            color: #fff;
            box-shadow: 0 22px 55px rgba(86, 0, 74, 0.15);
            margin-bottom: 18px;
        }

        .case-hero::after {
            content: "";
            position: absolute;
            right: -20px;
            bottom: -44px;
            width: 260px;
            height: 150px;
            border: 1px solid rgba(255,255,255,0.28);
            border-radius: 50%;
            transform: rotate(-15deg);
        }

        .brand-logo {
            position: absolute;
            top: 24px;
            right: 28px;
            z-index: 1;
            width: 240px;
            max-width: 28%;
            padding: 14px 16px 12px;
            border: 1px solid rgba(255,255,255,0.55);
            border-radius: 8px;
            background: rgba(255,255,255,0.92);
            box-shadow: 0 14px 28px rgba(86, 0, 74, 0.16);
        }

        .brand-logo svg {
            display: block;
            width: 100%;
            height: auto;
        }

        .case-kicker {
            position: relative;
            display: inline-flex;
            align-items: center;
            gap: 8px;
            padding: 5px 10px;
            border: 1px solid rgba(255,255,255,0.22);
            border-radius: 6px;
            background: rgba(255,255,255,0.14);
            font-size: 13px;
            font-weight: 700;
        }

        .case-hero h1 {
            position: relative;
            max-width: 780px;
            margin: 16px 0 10px;
            color: #fff;
            font-size: clamp(30px, 4.2vw, 48px);
            line-height: 1.08;
            font-weight: 850;
        }

        .case-hero p {
            position: relative;
            max-width: 800px;
            margin: 0;
            color: rgba(255,255,255,0.86);
            font-size: 16px;
            line-height: 1.75;
        }

        .metric-grid {
            display: grid;
            grid-template-columns: repeat(4, minmax(0, 1fr));
            gap: 12px;
            margin: 18px 0 24px;
        }

        .metric-card {
            min-height: 86px;
            border: 1px solid var(--line);
            border-radius: 8px;
            padding: 14px 15px;
            background: rgba(255,255,255,0.88);
            box-shadow: 0 12px 28px rgba(70, 58, 76, 0.06);
        }

        .metric-label {
            margin: 0 0 5px;
            color: var(--muted);
            font-size: 12px;
            font-weight: 700;
        }

        .metric-value {
            margin: 0;
            color: var(--ink);
            font-size: 18px;
            font-weight: 800;
        }

        .section-card {
            border: 1px solid var(--line);
            border-left: 5px solid var(--primary);
            border-radius: 8px;
            padding: 18px 20px;
            margin: 16px 0 14px;
            background: var(--panel);
            box-shadow: var(--shadow);
        }

        .section-card.accent {
            border-left-color: var(--accent);
        }

        .section-kicker {
            margin: 0 0 6px;
            color: var(--primary);
            font-size: 12px;
            font-weight: 800;
        }

        .section-card.accent .section-kicker {
            color: var(--accent);
        }

        .section-title {
            margin: 0 0 6px;
            color: var(--ink);
            font-size: 22px;
            font-weight: 850;
        }

        .section-copy {
            margin: 0;
            color: var(--muted);
            font-size: 14px;
            line-height: 1.7;
        }

        .guideline-panel {
            border: 1px solid var(--warning-line);
            border-radius: 8px;
            padding: 15px 17px;
            margin: 4px 0 18px;
            background: var(--warning-bg);
            color: #6b3154;
            line-height: 1.75;
        }

        .footer-note {
            margin-top: 24px;
            padding: 13px 16px;
            border-top: 1px solid var(--line);
            color: var(--muted);
            font-size: 13px;
            text-align: center;
        }

        @media (max-width: 900px) {
            .metric-grid {
                grid-template-columns: repeat(2, minmax(0, 1fr));
            }

            .case-hero {
                padding: 28px 22px;
            }

            .brand-logo {
                position: relative;
                top: auto;
                right: auto;
                width: 210px;
                max-width: 100%;
                margin-bottom: 22px;
            }
        }

        @media (max-width: 620px) {
            .metric-grid {
                grid-template-columns: 1fr;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_workbench_header():
    st.markdown(
        """
        <section class="case-hero">
          <div class="brand-logo" aria-label="和黄医药 HUTCHMED logo">
            <svg viewBox="0 0 400 190" role="img" aria-labelledby="logo-title">
              <title id="logo-title">和黄医药 HUTCHMED</title>
              <g fill="none" stroke-linecap="round">
                <path d="M210 12v54" stroke="#b00086" stroke-width="18"/>
                <path d="M244 44v25" stroke="#b00086" stroke-width="16"/>
                <path d="M276 41v31" stroke="#ec008c" stroke-width="16"/>
                <path d="M308 43v31" stroke="#ec008c" stroke-width="16"/>
                <path d="M340 45v29" stroke="#ec008c" stroke-width="16"/>
                <path d="M364 47v44" stroke="#ec008c" stroke-width="16"/>
                <path d="M386 52v22" stroke="#ec008c" stroke-width="16"/>
              </g>
              <text x="0" y="139" fill="#b00086" font-size="62" font-weight="900" font-family="Microsoft YaHei, SimHei, sans-serif">和黄</text>
              <text x="150" y="139" fill="#ec008c" font-size="62" font-weight="900" font-family="Microsoft YaHei, SimHei, sans-serif">医药</text>
              <text x="2" y="178" fill="#b00086" font-size="31" font-weight="900" letter-spacing="0" font-family="Arial, sans-serif">HUTCHMED</text>
            </svg>
          </div>
          <div class="case-kicker">Case Slide Studio · Oncology Workflow</div>
          <h1>肿瘤病例 PPT 自动生成工作台</h1>
          <p>将病历截图或电子病历转换为结构化病例资料，并自动生成适合病例讨论、科室汇报和学术交流的 PowerPoint。</p>
        </section>
        <section class="metric-grid" aria-label="功能概览">
          <article class="metric-card">
            <p class="metric-label">入口</p>
            <p class="metric-value">OCR / 文本</p>
          </article>
          <article class="metric-card">
            <p class="metric-label">结构化</p>
            <p class="metric-value">分线推断</p>
          </article>
          <article class="metric-card">
            <p class="metric-label">输出</p>
            <p class="metric-value">PPTX 幻灯片</p>
          </article>
          <article class="metric-card">
            <p class="metric-label">附加</p>
            <p class="metric-value">逻辑流摘要</p>
          </article>
        </section>
        """,
        unsafe_allow_html=True,
    )


def render_section_card(kicker, title, copy, accent=False):
    accent_class = " accent" if accent else ""
    st.markdown(
        f"""
        <article class="section-card{accent_class}">
          <p class="section-kicker">{kicker}</p>
          <h2 class="section-title">{title}</h2>
          <p class="section-copy">{copy}</p>
        </article>
        """,
        unsafe_allow_html=True,
    )


def render_generated_case(case_json, ppt_file, file_name):
    st.success("深度解析成功！您可以下载完整 PPT，或直接复制下方的逻辑流。")

    col1, col2 = st.columns([2, 1])
    with col1:
        st.download_button(
            label="立即下载完整 PPT",
            data=ppt_file,
            file_name=file_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    with col2:
        with st.expander("查看底层 JSON 树"):
            st.json(case_json)

    st.markdown("### 病例全病程逻辑线")
    st.info(render_logic_line_markdown(case_json))


inject_frontend_style()
render_workbench_header()

tab1, tab2 = st.tabs(["上传图片识别", "粘贴电子病历"])

if "ocr_result_text" not in st.session_state:
    st.session_state.ocr_result_text = ""

with tab1:
    render_section_card(
        "STEP 01",
        "批量上传病历图片",
        "适合处理病历截图、出院小结、检查报告和影像结论截图。建议按病程顺序上传，便于后续时间轴整理。",
    )

    st.markdown(
        """
        <div class="guideline-panel">
          <strong>上传建议</strong>
          <ul>
            <li>优先上传电脑原图截图或高清扫描件，推荐使用微信 Alt+A 截图后保存。</li>
            <li>如果拍摄电脑屏幕，请保持界面清晰，减少摩尔纹和反光干扰。</li>
            <li>单张图片建议控制在 4MB 以内，系统会对过大图片做轻度压缩。</li>
          </ul>
        </div>
        """,
        unsafe_allow_html=True,
    )

    uploaded_files = st.file_uploader(
        "上传病历截图、报告图片或出院小结",
        type=["png", "jpg", "jpeg"],
        accept_multiple_files=True,
    )

    if uploaded_files:
        st.info(f"已选择 {len(uploaded_files)} 张图片。")
        if st.button("开始批量提取文字"):
            with st.spinner("正在呼叫百度高精度 OCR 引擎扫描所有图片..."):
                token = get_baidu_access_token()
                if not token:
                    st.error("获取百度 API 授权失败，请检查密钥。")
                else:
                    all_extracted_text = []
                    for i, file in enumerate(uploaded_files):
                        image_bytes = file.getvalue()
                        text = perform_ocr(image_bytes, token)
                        all_extracted_text.append(f"【第 {i + 1} 页提取结果】\n{text}\n")
                    st.session_state.ocr_result_text = "\n".join(all_extracted_text)
            st.success("文字提取成功！请在下方核对。")

    render_section_card(
        "STEP 02",
        "人工校对与补全",
        "请补充 OCR 未识别到的关键诊疗信息，包括病理、基因检测、治疗方案、疗效评估和本次入院计划。",
        accent=True,
    )

    final_text_to_process = st.text_area(
        "校对并补全病史",
        value=st.session_state.ocr_result_text,
        height=350,
        placeholder="在这里核对 OCR 结果，也可以直接补充病史、检查、治疗和疗效信息。",
    )

    if st.button("自动推断分线并生成 PPT", type="primary"):
        if len(final_text_to_process) < 20:
            st.warning("病史太短，请补充详细记录。")
        else:
            try:
                with st.spinner("AI 正在按时间轴拆解并自动推断治疗线数..."):
                    case_json = extract_complex_case(final_text_to_process)
                with st.spinner("正在绘制时间轴并排版幻灯片..."):
                    maker = AdvancedPPTMaker(case_json)
                    ppt_file = maker.build()

                render_generated_case(case_json, ppt_file, "病例汇报_OCR版.pptx")

            except Exception as e:
                st.error(f"运行出错，请核对：{str(e)}")

with tab2:
    render_section_card(
        "TEXT MODE",
        "粘贴电子病历",
        "适合直接从医院系统、Word 文档或随访记录中复制长病历。内容越完整，PPT 分线、时间轴和讨论点越稳定。",
        accent=True,
    )

    patient_input = st.text_area(
        "请贴入详细病史",
        height=300,
        placeholder="建议包含：基本信息、主诉、诊断、病理/分子检测、每线治疗方案、疗效评估、肿瘤标志物变化、本次入院资料和后续计划。",
    )

    if st.button("开始深度解析并生成 PPT", key="btn_text", type="primary"):
        if len(patient_input) < 20:
            st.warning("病史太短，请提供详细病历。")
        else:
            try:
                with st.spinner("AI 正在按时间轴拆解并自动推断治疗线数..."):
                    case_json = extract_complex_case(patient_input)
                with st.spinner("正在自动排版幻灯片..."):
                    maker = AdvancedPPTMaker(case_json)
                    ppt_file = maker.build()

                render_generated_case(case_json, ppt_file, "病例汇报_文本版.pptx")

            except Exception as e:
                st.error(f"运行出错，请核对：{str(e)}")

st.markdown(
    '<div class="footer-note">页面已接入 OCR、AI 解析与 PPT 生成流程；请在部署环境配置 Streamlit Secrets 后使用。</div>',
    unsafe_allow_html=True,
)
